import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Add background shape
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(11, 17, 32)
bg.line.fill.background()

# Title text
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "GeoConvey Presentation Test"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(16, 185, 129)

prs.save("test.pptx")
print("Test presentation saved successfully")
